from pptx import Presentation
from itertools import groupby
import json
from io import BytesIO

list_of_notes = []
list_of_notes_address = []

class parser: 
    def __init__(self):
        self.ppt = None

    def load(self, file) -> None:
        self.ppt = Presentation(file)
    
    #Получение массива адресов страниц
    def update(ppt):
        #Получение адресов закладок/страниц
        for slide in enumerate(ppt.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        note_address = run.hyperlink.address
                        if (note_address != None):
                            list_of_notes_address.append(note_address)
                        else:
                            continue
        #Удаление повторяющихся адресов, возникающих из-за пробелов и других символов
        new_list_of_notes_address = [el for el, _ in groupby(list_of_notes_address)]

        notes = parser.get_text_from_slides(ppt)
        parser.merging_lists(new_list_of_notes_address, notes, list_of_notes)

    #Получение текста со слайдов
    def get_text_from_slides(ppt):
        notes = []
        for page, slide in enumerate(ppt.slides):
            temp = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    temp.append(shape.text)       
            temp = parser.check_slash_n(temp)
            notes.append(temp)
        return notes

    #Объединение массива адресов и массива текстов в один
    def merging_lists(list1, list2, list3):
        for i in range(0, max(len(list1), len(list2))):
            if i < len(list1):
                list3.append(list1[i])
            if i < len(list2):
                list3.append(list2[i])
        print(list3)
        return list3

    #Удаление символов табуляции и перехода на другую строку из массива
    def check_slash_n(list):
        list = [x.replace('\n', ' ') for x in list]
        list = [x.replace('\x0b', ' ') for x in list]
        list = [x.replace('\r', ' ') for x in list]
        list = [x.replace('\t', ' ') for x in list]
        return list
    
    #Сохранение массива в json файл
    def save_to_json():
        with open('notes.json', 'w', encoding='utf-8') as file:
            json.dump(list_of_notes, file, ensure_ascii=False, indent=4)

    #Строка для ввода названия презентации
    def get_pres_name():
        print('Enter the name of your presentation: ')

    def save_to_bytes(self) -> BytesIO:
        if not self.ppt:
            raise ValueError("pptx файл не загружен.")
        output = BytesIO()
        self.workbook.save(output)
        output.seek(0)
        return output

#Ввод имени презентации(без ".pptx")
#parser.get_pres_name()
#presentation_name = input()
#presentation = Presentation(presentation_name + '.pptx')

#Парсинг презентации и сохранение в json формат
#parser.update(presentation)
#parser.save_to_json()
