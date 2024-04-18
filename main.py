from pptx import Presentation
from itertools import groupby
import json

list_of_notes = [] #Массив закладок
list_of_notes_address = [] #Массив адресов закладок

class parser:
    #Получение массива адресов страниц
    def get_notes(ppt):
        #Получение адресов закладок/страниц
        for slide in ppt.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        note_address = run.hyperlink.address
                        if (note_address != None):
                            list_of_notes_address.append(note_address)
                        else:
                            continue
        #Удаление повторяющихся адресов, возникающих из-за пробелов и других символов
        new_list_of_notes_address = [el for el, _ in groupby(list_of_notes_address)]

        #Получение текста со слайдов
        notes = []
        for page, slide in enumerate(ppt.slides):
            temp = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    temp.append(shape.text) 
            #Удаление символов табуляции и перехода на другую строку из массива      
            temp = [x.replace('\n', ' ') for x in temp]
            temp = [x.replace('\x0b', ' ') for x in temp]
            temp = [x.replace('\r', ' ') for x in temp]
            temp = [x.replace('\t', ' ') for x in temp]
            notes.append(temp)
        print(notes)

        #Объединение массива адресов и массива текстов в один
        for i in range(0, max(len(new_list_of_notes_address), len(notes))):
            if i < len(new_list_of_notes_address):
                list_of_notes.append(new_list_of_notes_address[i])
            if i < len(notes):
                list_of_notes.append(notes[i])
        return list_of_notes

    #Сохранение массива в json файл
    def save_to_json():
        with open('notes.json', 'w', encoding='utf-8') as file:
            json.dump(list_of_notes, file, ensure_ascii=False, indent=4)

    #Строка для ввода названия презентации
    def get_pres_name():
        print('Enter the name of your presentation: ')

#Ввод имени презентации(без ".pptx")
parser.get_pres_name()
presentation_name = input()
presentation = Presentation(presentation_name + '.pptx')

#Парсинг презентации и сохранение в json формат
parser.get_notes(presentation)
parser.save_to_json()

